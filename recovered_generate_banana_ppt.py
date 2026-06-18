 """
 generate_banana_ppt.py
 ======================
 Generates a highly-styled, modern, and visually premium 11-slide PowerPoint
 presentation for the Four-Group BFT Consensus Comparison, following the Canva
 "Studio Shodwe 2024 Project" layout structure and the "banana" aesthetic.
 """
 
 import os
 import sys
 
 try:
     from pptx import Presentation
     from pptx.util import Inches, Pt
     from pptx.dml.color import RGBColor
     from pptx.enum.text import PP_ALIGN
     from pptx.enum.shapes import MSO_SHAPE
     from pptx.oxml.xmlchemy import OxmlElement
 except ImportError:
     print("Installing python-pptx...")
     os.system(f"{sys.executable} -m pip install python-pptx")
     from pptx import Presentation
     from pptx.util import Inches, Pt
     from pptx.dml.color import RGBColor
     from pptx.enum.text import PP_ALIGN
     from pptx.enum.shapes import MSO_SHAPE
     from pptx.oxml.xmlchemy import OxmlElement
 
 # ── Paths ────────────────────────────────────────────────────────────────────
 BASE_DIR = os.path.dirname(os.path.abspath(__file__))
 FIGURE_DIR = os.path.join(BASE_DIR, "figures")
 BACKGROUND_IMG = os.path.join(BASE_DIR, "slide_background_generated.png")
 OUTPUT_FILE = os.path.join(BASE_DIR, "four_group_comparison_presentation_v8_banana.pptx")
 
 # ── Color Palette ──────
<truncated 33077 bytes>
 ]
     add_bullet_text(slide, bullets_s11_l, 0.7, 1.7, 5.4, 4.9, font_size=13)
 
     add_card(slide, 6.6, 1.5, 6.2, 5.3)
     add_title_text(slide, "Key Academic References", 6.8, 1.7, 5.8, 0.4, font_size=16, color=TEXT_DARK)
     refs = [
         "[1] A. Sheikh et al., 'Secured Energy Trading Using Byzantine-Based Blockchain Consensus,' IEEE Access, 2020.",
         "[2] L. Lamport et al., 'The Byzantine Generals Problem,' ACM TOPLAS, 1982.",
         "[3] X. Ding et al., 'CE-PBFT: Optimized PBFT Consensus for Microgrid Power Trading,' IEEE Trans. Smart Grid, 2024.",
         "[4] Z. Xu et al., 'G-PBFT Algorithm and its Application in Distributed Energy Trading,' Applied Energy, 2025.",
         "[5] Y. Suo et al., 'SV-PBFT: Consensus Algorithm for V2V Energy Transactions,' IEEE IoT Journal, 2025.",
         "[6] H. Wang et al., 'RVR Blockchain Consensus for Smart Grid Energy Trading,' IEEE TPDS, 2025."
     ]
     add_bullet_text(slide, refs, 6.8, 2.1, 5.8, 4.5, font_size=11, color=TEXT_MID, bold_first=False)
 
     add_speaker_notes(slide,
         "We conclude with Slide 11: Thank You. "
         "Consensus latency is the critical parameter governing physical smart grid security. "
         "We recommend Group 4 protocols for modern grids and Group 3 for legacy grids. "
         "I would like to thank our mentors, Uday sir and Sunny sir, and I am now open to your questions."
     )
 
     # ── Save ─────────────────────────────────────────────────────────────
     prs.save(OUTPUT_FILE)
     print(f"\n  [OK] Banana Presentation saved: {OUTPUT_FILE}")
     print(f"  [OK] Total slides: {len(prs.slides)}")
     return OUTPUT_FILE
 
 if __name__ == "__main__":
     create_presentation()
 
The above content shows the entire, complete file contents of the requested file.
