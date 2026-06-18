 """
 generate_premium_ppt.py
 ==========================
 Generates a highly-styled, modern PowerPoint presentation for the
 Four-Group BFT Consensus Comparison.
 """
 
 import os
 import sys
 
 try:
     from pptx import Presentation
     from pptx.util import Inches, Pt
     from pptx.dml.color import RGBColor
     from pptx.enum.text import PP_ALIGN
     from pptx.enum.shapes import MSO_SHAPE
     from pptx.oxml.xmlchemy import OxmlElement
 except ImportError:
     print("Installing python-pptx...")
     os.system(f"{sys.executable} -m pip install python-pptx")
     from pptx import Presentation
     from pptx.util import Inches, Pt
     from pptx.dml.color import RGBColor
     from pptx.enum.text import PP_ALIGN
     from pptx.enum.shapes import MSO_SHAPE
     from pptx.oxml.xmlchemy import OxmlElement
 
 # ── Paths ────────────────────────────────────────────────────────────────────
 BASE_DIR = os.path.dirname(os.path.abspath(__file__))
 FIGURE_DIR = os.path.join(BASE_DIR, "figures")
 BACKGROUND_IMG = os.path.join(BASE_DIR, "slide_background_generated.png")
 OUTPUT_FILE = os.path.join(BASE_DIR, "four_group_comparison_presentation_v7_premium.pptx")
 
 # ── Color palette ───────────────────────────────────────────────
<truncated 26246 bytes>
out sub-ms sync"
         ])
     ]
     add_bullet_text(slide, summary_bullets, 0.7, 1.6, 11.9, 5.0, font_size=16)
 
     # ════════ SLIDE 16: References ════════
     slide = prs.slides.add_slide(prs.slide_layouts[6])
     add_slide_background(slide, prs)
     add_title_text(slide, "15. References", 0.5, 0.3, 12.0, 0.7, font_size=30, color=ACCENT)
     add_accent_bar(slide, 0.5, 1.1, 2.5, 0.08)
 
     add_card(slide, 0.5, 1.4, 12.3, 5.5)
     refs = [
         "[1] A. Sheikh et al., 'Secured Energy Trading Using Byzantine-Based Blockchain Consensus,' IEEE Access, vol. 8, 2020.",
         "[2] L. Lamport et al., 'The Byzantine Generals Problem,' ACM TOPLAS, vol. 4, no. 3, 1982.",
         "[3] M. Castro and B. Liskov, 'Practical Byzantine Fault Tolerance,' in Proc. OSDI, 1999.",
         "[4] X. Ding et al., 'CE-PBFT: An Optimized PBFT Consensus Algorithm for Microgrid Power Trading,' IEEE Trans. Smart Grid, 2024.",
         "[5] Z. Xu et al., 'G-PBFT Algorithm and its Application in Distributed Energy Trading,' Applied Energy, 2025.",
         "[6] Y. Suo et al., 'SV-PBFT: An Efficient and Stable Blockchain PBFT Improved Consensus Algorithm for Vehicle-to-Vehicle Energy Transactions,' IEEE IoT Journal, 2025.",
         "[7] A. Yakovenko, 'Solana: A new architecture for a high performance blockchain,' Solana Whitepaper, 2018.",
         "[8] H. Wang et al., 'RVR Consensus: A Verifiable, Weighted-Random, Byzantine-Tolerant Smart Grid Framework,' IEEE TPDS, 2025."
     ]
     add_bullet_text(slide, refs, 0.7, 1.6, 11.9, 5.0, font_size=13, color=TEXT_MID, bold_first=False)
 
     # ════════ SAVE ════════
     prs.save(OUTPUT_FILE)
     print(f"\n  [OK] Premium Presentation saved: {OUTPUT_FILE}")
 
 if __name__ == "__main__":
     create_presentation()
 
The above content shows the entire, complete file contents of the requested file.
