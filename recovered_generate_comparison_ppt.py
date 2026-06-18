 """
 generate_comparison_ppt.py
 ==========================
 Generates a professional PowerPoint presentation for the
 Four-Group BFT Consensus Comparison, following the official VJTI guidelines.
 
 Usage:
     python generate_comparison_ppt.py
 
 Output:
     ./four_group_comparison_presentation_updated.pptx
 """
 
 import os
 import sys
 
 try:
     from pptx import Presentation
     from pptx.util import Inches, Pt
     from pptx.dml.color import RGBColor
     from pptx.enum.text import PP_ALIGN
     from pptx.enum.shapes import MSO_SHAPE
 except ImportError:
     print("Installing python-pptx...")
     os.system(f"{sys.executable} -m pip install python-pptx")
     from pptx import Presentation
     from pptx.util import Inches, Pt
     from pptx.dml.color import RGBColor
     from pptx.enum.text import PP_ALIGN
     from pptx.enum.shapes import MSO_SHAPE
 
 # ── Paths ────────────────────────────────────────────────────────────────────
 BASE_DIR = os.path.dirname(os.path.abspath(__file__))
 FIGURE_DIR = os.path.join(BASE_DIR, "figures")
 BACKGROUND_IMG = os.path.join(BASE_DIR, "slide_background_generated.png")
 OUTPUT_FILE = os.path.join(BASE_DIR, "four_group_comparison_presentation_v4.pptx")
 
 # ── Color palette ───────────────
<truncated 43057 bytes>
mport et al., 'The Byzantine Generals Problem,' ACM TOPLAS, vol. 4, no. 3, 1982.",
         "[3] M. Castro and B. Liskov, 'Practical Byzantine Fault Tolerance,' in Proc. OSDI, 1999.",
         "[4] X. Ding et al., 'CE-PBFT: An Optimized PBFT Consensus Algorithm for Microgrid Power Trading,' IEEE Trans. Smart Grid, 2024.",
         "[5] Z. Xu et al., 'G-PBFT Algorithm and its Application in Distributed Energy Trading,' Applied Energy, 2025.",
         "[6] Y. Suo et al., 'SV-PBFT: An Efficient and Stable Blockchain PBFT Improved Consensus Algorithm for Vehicle-to-Vehicle Energy Transactions,' IEEE IoT Journal, 2025.",
         "[7] A. Yakovenko, 'Solana: A new architecture for a high performance blockchain,' Solana Whitepaper, 2018.",
         "[8] H. Wang et al., 'RVR Consensus: A Verifiable, Weighted-Random, Byzantine-Tolerant Smart Grid Framework,' IEEE TPDS, 2025."
     ]
     add_bullet_text(slide, refs, 0.5, 1.3, 12.3, 5.2, font_size=12, color=TEXT_MID)
     
     add_speaker_notes(slide,
         "Here is the list of core references from the literature. These form the baseline and comparison algorithms we evaluated. "
         "Thank you for your attention, and I am now open to any questions or suggestions from the mentors."
     )
 
     # ── Save ─────────────────────────────────────────────────────────────
     prs.save(OUTPUT_FILE)
     print(f"\n  [OK] Presentation saved: {OUTPUT_FILE}")
     print(f"  [OK] Total slides: {len(prs.slides)}")
     return OUTPUT_FILE
 
 
 if __name__ == "__main__":
     print("=" * 60)
     print("  GENERATING COMPLIANT FOUR-GROUP BFT COMPARISON PRESENTATION")
     print("=" * 60)
     create_presentation()
     print("=" * 60)
 
The above content shows the entire, complete file contents of the requested file.
